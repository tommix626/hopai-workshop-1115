from pptx import Presentation #pip install python-pptx

#pip install PyPDF2
#pip install langchain openai

def extract_text_from_ppt(ppt_path):
    presentation = Presentation(ppt_path)
    text_data = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_data.append(shape.text)

    return text_data

if __name__ == '__main__':


    ppt_text = extract_text_from_ppt('test.pptx')
    print(ppt_text)

# See PyCharm help at https://www.jetbrains.com/help/pycharm/
